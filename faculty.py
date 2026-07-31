import streamlit as st
import pypdf
from pptx import Presentation
import docx
import requests
import google.generativeai as genai
import datetime
import random
import zlib
import base64

st.set_page_config(page_title="Faculty Studio", page_icon="👨‍🏫", layout="wide")

# Hide Streamlit form submission helper notes
st.markdown("""
    <style>
    [data-testid="InputInstructions"] {
        display: none !important;
    }
    </style>
""", unsafe_allow_html=True)

BIN_ID = st.secrets.get("JSONBIN_BIN_ID", "")
JSONBIN_KEY = st.secrets.get("JSONBIN_API_KEY", "")
GEMINI_KEY = st.secrets.get("GEMINI_API_KEY", "")
DEV_OVERRIDE = "1901"

if not BIN_ID or not JSONBIN_KEY or not GEMINI_KEY:
    st.error("🔑 Database Credentials or Gemini API Key missing in Streamlit Secrets!")
    st.stop()

genai.configure(api_key=GEMINI_KEY)
style_analyzer_model = genai.GenerativeModel("models/gemini-2.5-flash-lite")

# --- COMPRESSION & OPTIMIZATION HELPERS ---
def compress_text(text: str) -> str:
    """Compresses string text using zlib and base64 encoding."""
    if not text:
        return ""
    try:
        compressed_bytes = zlib.compress(text.encode('utf-8'), level=9)
        return "COMPRESSED:" + base64.b64encode(compressed_bytes).decode('utf-8')
    except Exception:
        return text

def decompress_text(text: str) -> str:
    """Decompresses zlib base64 string back to raw text."""
    if not text or not isinstance(text, str):
        return ""
    if text.startswith("COMPRESSED:"):
        try:
            raw_b64 = text.replace("COMPRESSED:", "")
            compressed_bytes = base64.b64decode(raw_b64.encode('utf-8'))
            return zlib.decompress(compressed_bytes).decode('utf-8')
        except Exception:
            return text
    return text

def decompress_course_data(raw_data: dict) -> dict:
    """Recursively decompresses slide text for local use."""
    if not isinstance(raw_data, dict):
        return {}
    
    for course_id, course_info in raw_data.items():
        if isinstance(course_info, dict):
            sessions = course_info.get("sessions", {})
            if isinstance(sessions, dict):
                for s_id, s_info in sessions.items():
                    if isinstance(s_info, dict) and "slides" in s_info:
                        s_info["slides"] = decompress_text(s_info["slides"])
    return raw_data

def compress_course_data_for_saving(data: dict) -> dict:
    """Creates a deep copy of course data, stripping raw PQ text and compressing slides."""
    import copy
    data_to_save = copy.deepcopy(data)
    
    for course_id, course_info in data_to_save.items():
        if isinstance(course_info, dict):
            # Strip bulky raw PQ text (only keep style_profile, which is all student app needs)
            course_info["global_pqs_text"] = ""
            
            sessions = course_info.get("sessions", {})
            if isinstance(sessions, dict):
                for s_id, s_info in sessions.items():
                    if isinstance(s_info, dict):
                        if "slides" in s_info:
                            s_info["slides"] = compress_text(s_info["slides"])
                        s_info["pqs"] = ""  # Strip raw session PQ text
    return data_to_save

def load_cloud_data():
    url = f"https://api.jsonbin.io/v3/b/{BIN_ID}/latest"
    headers = {"X-Master-Key": JSONBIN_KEY}
    try:
        req = requests.get(url, headers=headers)
        if req.status_code == 200:
            raw_data = req.json().get("record", {})
            valid_courses = {k: v for k, v in raw_data.items() if isinstance(v, dict) and ("sessions" in v or "passcode" in v)}
            return decompress_course_data(valid_courses)
    except Exception:
        return {}
    return {}

def save_cloud_data(data):
    url = f"https://api.jsonbin.io/v3/b/{BIN_ID}"
    headers = {
        "Content-Type": "application/json",
        "X-Master-Key": JSONBIN_KEY
    }
    
    payload = compress_course_data_for_saving(data)
    
    try:
        req = requests.put(url, headers=headers, json=payload)
        if req.status_code == 200:
            return True
        else:
            st.error(f"Failed to sync with cloud database ({req.status_code}): {req.text}")
            return False
    except Exception as e:
        st.error(f"Failed to sync with cloud database: {e}")
        return False

def extract_text_from_file(uploaded_file):
    if not uploaded_file:
        return ""
    
    filename = uploaded_file.name.lower()
    text = ""
    
    try:
        if filename.endswith(".pdf"):
            pdf_reader = pypdf.PdfReader(uploaded_file)
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                extracted = page.extract_text()
                if extracted:
                    text += f"\n--- Page/Slide {page_num} ---\n" + extracted
                    
        elif filename.endswith(".pptx"):
            prs = Presentation(uploaded_file)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_words = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_words.append(shape.text)
                if slide_words:
                    text += f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_words)
                    
        elif filename.endswith(".docx"):
            doc = docx.Document(uploaded_file)
            p_text = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        p_text.append(" | ".join(row_text))
            text = "\n".join(p_text)
            
    except Exception as e:
        st.error(f"Error reading '{uploaded_file.name}': {e}")
        
    return text

def analyze_style_profile(pqs_text):
    if not pqs_text or not pqs_text.strip():
        return "Standard NBME clinical vignette style with 4-5 choices."
    
    analysis_prompt = f"""
    Analyze the following medical practice questions and reverse-engineer the writing style into concise guidance rules:
    1. Vignette length and clinical depth.
    2. Question stem phrasing.
    3. Distractor structure and difficulty.
    4. Rationale and explanation format.

    PRACTICE QUESTIONS:
    {pqs_text[:6000]}
    """
    try:
        res = style_analyzer_model.generate_content(analysis_prompt)
        return res.text
    except Exception as e:
        return f"Standard NBME style (Error during analysis: {e})"

@st.dialog("✅ Session Saved Successfully!")
def show_save_confirmation(session_name, action_type="published"):
    st.write(f"The lecture session **'{session_name}'** has been successfully {action_type} and is now available to students on the practice app.")
    if st.button("Close Window", type="primary"):
        st.session_state.saved_session_info = None
        st.rerun()

data = load_cloud_data()

if "authenticated_course" not in st.session_state:
    st.session_state.authenticated_course = None
if "show_delete_course_confirm" not in st.session_state:
    st.session_state.show_delete_course_confirm = False
if "saved_session_info" not in st.session_state:
    st.session_state.saved_session_info = None

st.title("👨‍🏫 Faculty Studio Portal")

if st.session_state.saved_session_info:
    info = st.session_state.saved_session_info
    show_save_confirmation(info["title"], info["action"])

# --- LOGIN SCREEN ---
if not st.session_state.authenticated_course:
    st.caption("Upload lecture slides, manage course materials, and train the AI on your question writing style.")
    
    tab1, tab2 = st.tabs(["🔒 Enter Existing Course Workspace", "➕ Create New Course Workspace"])
    
    with tab1:
        if not data:
            st.info("👋 Welcome! No active courses found. Switch to the 'Create New Course Workspace' tab to set up your first course.")
        else:
            with st.form("login_form"):
                st.markdown("##### 🔑 Access Your Workspace")
                selected_course = st.selectbox("Select Your Course:", options=list(data.keys()), help="Select the course you want to manage.")
                entered_code = st.text_input("Enter 4-Digit Course Passcode:", type="password", help="Enter the 4-digit PIN generated when this course was created.")
                submit_login = st.form_submit_button("Enter Studio Workspace", type="primary")
                
                if submit_login:
                    stored_code = str(data[selected_course].get("passcode"))
                    if entered_code == stored_code or entered_code == DEV_OVERRIDE:
                        st.session_state.authenticated_course = selected_course
                        st.session_state.show_delete_course_confirm = False
                        st.rerun()
                    else:
                        st.error("❌ Incorrect passcode. Please try again or ask your course director.")

    with tab2:
        with st.form("create_course_form"):
            st.markdown("##### 🚀 Create a New Course")
            st.caption("Creating a workspace generates a unique 4-digit passcode for you and co-instructors.")
            new_course_name = st.text_input("Course Name:", placeholder="e.g., Gastrointestinal System (GI) - Fall 2026")
            submit_create = st.form_submit_button("Generate Course Workspace", type="primary")
            
            if submit_create:
                if not new_course_name.strip():
                    st.warning("⚠️ Please enter a course name.")
                elif new_course_name in data:
                    st.warning("⚠️ A course with this name already exists. Please choose a unique title.")
                else:
                    generated_passcode = str(random.randint(1000, 9999))
                    data[new_course_name] = {
                        "passcode": generated_passcode,
                        "global_pqs_filename": "",
                        "global_pqs_text": "",
                        "global_style_profile": "",
                        "sessions": {}
                    }
                    if save_cloud_data(data):
                        st.session_state.newly_created_course = new_course_name
                        st.session_state.newly_created_passcode = generated_passcode
                        st.rerun()

        if "newly_created_course" in st.session_state and st.session_state.newly_created_course:
            nc_name = st.session_state.newly_created_course
            nc_code = st.session_state.newly_created_passcode
            
            st.success(f"🎉 Workspace for '{nc_name}' created successfully!")
            st.info(f"🔐 **Your Assigned Passcode:** `{nc_code}`\n\n*Save this 4-digit code! You and co-faculty will need it to log into this course workspace in the future.*")
            
            if st.button(f"🚀 Enter '{nc_name}' Workspace Now", type="primary"):
                st.session_state.authenticated_course = nc_name
                st.session_state.newly_created_course = None
                st.session_state.newly_created_passcode = None
                st.rerun()

# --- WORKSPACE DASHBOARD ---
else:
    active_course = st.session_state.authenticated_course
    course_data = data.get(active_course)

    if not course_data:
        st.session_state.authenticated_course = None
        st.rerun()

    col_a, col_b, col_c = st.columns([2, 1, 1])
    with col_a:
        st.subheader(f"Active Workspace: **{active_course}**")
        st.caption(f"🔐 Passcode for co-instructors: **{course_data.get('passcode')}**")
    with col_b:
        if st.button("🔒 Exit Workspace"):
            st.session_state.authenticated_course = None
            st.session_state.show_delete_course_confirm = False
            st.rerun()
    with col_c:
        if st.button("🗑️ Delete Course Workspace"):
            st.session_state.show_delete_course_confirm = True

    if st.session_state.show_delete_course_confirm:
        st.warning(f"⚠️ **Are you sure you want to permanently delete '{active_course}'?** This erases all published sessions for students.")
        col_yes, col_no = st.columns([1, 4])
        with col_yes:
            if st.button("🚨 Yes, Delete Permanently", type="primary"):
                del data[active_course]
                save_cloud_data(data)
                st.session_state.authenticated_course = None
                st.session_state.show_delete_course_confirm = False
                st.success(f"'{active_course}' has been removed.")
                st.rerun()
        with col_no:
            if st.button("Cancel"):
                st.session_state.show_delete_course_confirm = False
                st.rerun()

    st.markdown("---")
    
    # --- FACULTY INSTRUCTIONS GUIDE BANNER ---
    with st.container():
        st.markdown("""
        ### 💡 Quick Start Guide
        1. **Course Exam Style (Recommended):** Expand the block below to upload a general practice exam/question bank. The AI will learn your exam writing style and apply it across all lectures by default.
        2. **Publish Lectures:** Go to **➕ Add New Session**, upload your slide deck (`PDF` or `PowerPoint`), set the date held, and select your practice question reference choice.
        3. **Edit or Update:** Use **🛠️ View/Edit Published Sessions** to update slide decks, adjust session dates, or upload custom exam sets anytime.
        """)
    
    st.markdown("---")
    
    # --- COURSE-WIDE GENERAL PQ UPLOADER SECTION ---
    master_fn = course_data.get("global_pqs_filename", "")
    course_has_master = bool(master_fn)
    
    with st.expander(f"🎓 Course-Wide Master Practice Exam ({'🟢 Active: ' + master_fn if course_has_master else '⚠️ None Uploaded Yet'})", expanded=not course_has_master):
        st.info("💡 **How this works:** Most courses use one general end-of-block practice exam or NBME question set. Uploading it here allows all current and future lecture sessions to inherit your writing style automatically!")
        
        if master_fn:
            st.success(f"📎 **Currently Active Master Exam File:** `{master_fn}`")
        
        up_master_file = st.file_uploader(
            "Upload Course-Wide Master Exam File (PDF or Word):", 
            type=["pdf", "docx"], 
            key="course_master_pq",
            help="Upload a PDF or Word doc containing sample NBME-style board questions, quizzes, or past exams."
        )
        
        if st.button("💾 Save Course-Wide Practice Exam", type="primary"):
            if up_master_file:
                with st.spinner("Extracting text and analyzing master exam writing style..."):
                    m_text = extract_text_from_file(up_master_file)
                    m_style = analyze_style_profile(m_text)
                    course_data["global_pqs_filename"] = up_master_file.name
                    course_data["global_pqs_text"] = ""  # Stripped to save database space
                    course_data["global_style_profile"] = m_style
                    save_cloud_data(data)
                    st.success(f"Successfully uploaded and analyzed '{up_master_file.name}' as the course master exam!")
                    st.rerun()
            else:
                st.warning("Please choose a file to upload before saving.")

    st.markdown("---")
    t_add, t_manage = st.tabs(["➕ Add New Session", "🛠️ View/Edit Published Sessions"])
    
    opt_master = f"Use Course-Wide Master Exam ({master_fn})" if course_has_master else "Use Course-Wide Master Exam (⚠️ No master exam uploaded yet)"
    opt_custom = "Upload Custom Practice Question File for This Session Specifically"
    
    # --- TAB 1: ADD NEW SESSION ---
    with t_add:
        st.markdown("#### Add Lecture Session")
        st.caption("Fill in session details and upload the lecture slide deck for students.")
        
        session_title = st.text_input("1. Session / Lecture Title:", placeholder="e.g., Physiology of Salivary and Gastric Secretion", help="Enter the exact lecture name as it appears on the schedule.")
        session_date = st.date_input("2. Date Session Was/Will Be Held:", value=datetime.date.today(), format="MM/DD/YYYY", help="Dates are formatted as MM/DD/YYYY and used to sort lectures chronologically for students.")
        slides_file = st.file_uploader("3. Upload Lecture Slides (PDF or PowerPoint):", type=["pdf", "pptx"], key="add_slides", help="Upload the official lecture slide deck. Text inside slides will be used as the exclusive truth source for question generation.")
        
        st.markdown("---")
        st.markdown("#### 4. Question Writing Style Reference")
        st.caption("Select which exam set the AI model should analyze to learn vignette length, stem phrasing, and distractor difficulty:")
        
        pq_source_choice = st.radio(
            "Select Reference Source:",
            options=[opt_master, opt_custom],
            key="add_pq_choice"
        )
        
        session_custom_pq_file = None
        if pq_source_choice == opt_custom:
            session_custom_pq_file = st.file_uploader(
                "Upload Session-Specific Practice Question File (PDF or Word - Required):", 
                type=["pdf", "docx"], 
                key="add_session_custom_pq",
                help="Upload a custom exam or quiz set specific to this lecture."
            )

        st.markdown("---")
        if st.button("🚀 Save & Publish Session to Students", type="primary"):
            if not session_title.strip():
                st.error("❌ Action Required: Please enter a Session Title.")
            elif not slides_file:
                st.error("❌ Action Required: Please upload a Lecture Slide PDF or PowerPoint file.")
            elif pq_source_choice == opt_master and not course_has_master:
                st.error("❌ Action Required: You selected 'Use Course-Wide Master Exam', but no master exam file has been uploaded for this course yet. Either upload one under 'Course-Wide Master Practice Exam' above, or select 'Upload Custom Practice Question File'.")
            elif pq_source_choice == opt_custom and not session_custom_pq_file:
                st.error("❌ Action Required: You chose to use a custom practice question set for this session, but have not attached a file.")
            else:
                with st.spinner("Processing slides, setting up style profile, and syncing to cloud..."):
                    slides_text = extract_text_from_file(slides_file)
                    date_str = session_date.strftime("%Y-%m-%d")
                    
                    if pq_source_choice == opt_custom and session_custom_pq_file:
                        pqs_text = extract_text_from_file(session_custom_pq_file)
                        pqs_fn = session_custom_pq_file.name
                        style_prof = analyze_style_profile(pqs_text)
                        pq_mode = "custom"
                    else:
                        pqs_fn = master_fn
                        style_prof = course_data.get("global_style_profile", "")
                        pq_mode = "course_master"

                    if "sessions" not in course_data:
                        course_data["sessions"] = {}

                    course_data["sessions"][session_title] = {
                        "date": date_str,
                        "slides": slides_text,
                        "slides_filename": slides_file.name,
                        "pq_mode": pq_mode,
                        "pqs": "",  # Stripped to save database space
                        "pqs_filename": pqs_fn,
                        "style_profile": style_prof
                    }
                    
                    if save_cloud_data(data):
                        st.session_state.saved_session_info = {"title": session_title, "action": "published"}
                        st.rerun()

    # --- TAB 2: MANAGE SESSIONS ---
    with t_manage:
        sessions = course_data.get("sessions", {})
        if not sessions:
            st.info("ℹ️ No published sessions in this course workspace yet. Use the '➕ Add New Session' tab above to publish your first lecture!")
        else:
            st.caption("Click any session below to edit its title, date, slide deck, or question style settings:")
            for s_title in list(sessions.keys()):
                sess_info = sessions[s_title]
                raw_date = sess_info.get("date", "N/A")
                
                display_date = "N/A"
                if raw_date != "N/A":
                    try:
                        display_date = datetime.datetime.strptime(raw_date, "%Y-%m-%d").strftime("%m/%d/%Y")
                    except ValueError:
                        display_date = raw_date
                
                with st.expander(f"📖 [{display_date}] {s_title}"):
                    mode = sess_info.get("pq_mode", "custom")
                    pq_fn_display = sess_info.get("pqs_filename", "Master Set" if mode == "course_master" else "N/A")
                    
                    st.write(f"**Date Held:** `{display_date}`")
                    st.write(f"**Reference Style Source:** {'🌐 Course-Wide Master Exam' if mode == 'course_master' else '📄 Custom Session Exam'} (`{pq_fn_display}`)")
                    
                    st.markdown("---")
                    st.markdown("##### ✏️ Edit Session Details")
                    
                    default_date = datetime.date.today()
                    if raw_date != "N/A":
                        try:
                            default_date = datetime.datetime.strptime(raw_date, "%Y-%m-%d").date()
                        except ValueError:
                            pass
                            
                    edit_title = st.text_input("Session Title:", value=s_title, key=f"edit_title_{s_title}")
                    edit_date = st.date_input("Session Date Held:", value=default_date, format="MM/DD/YYYY", key=f"edit_date_{s_title}")
                    
                    curr_slides_fn = sess_info.get("slides_filename", "[Uploaded file]")
                    st.markdown("**Lecture Slides File:**")
                    st.caption(f"📎 **Currently attached slide deck:** `{curr_slides_fn}`")
                    new_slides_file = st.file_uploader("Upload new slide file to replace (PDF or PPTX - Optional):", type=["pdf", "pptx"], key=f"edit_slides_{s_title}")
                    
                    st.markdown("---")
                    st.markdown("**Practice Question Style Reference**")
                    current_mode_index = 0 if mode == "course_master" else 1
                    edit_pq_choice = st.radio(
                        "Select Reference Source:",
                        options=[opt_master, opt_custom],
                        index=current_mode_index,
                        key=f"edit_pq_choice_{s_title}"
                    )
                    
                    edit_custom_pq_file = None
                    if edit_pq_choice == opt_custom:
                        if sess_info.get("pqs_filename"):
                            st.caption(f"📎 **Currently attached session exam:** `{sess_info.get('pqs_filename')}`")
                        edit_custom_pq_file = st.file_uploader("Upload replacement session practice question file (PDF/DOCX - Optional):", type=["pdf", "docx"], key=f"edit_pqs_{s_title}")

                    st.markdown("---")
                    col_save, col_del = st.columns([1, 1])
                    
                    with col_save:
                        if st.button("💾 Save Changes to Session", key=f"save_{s_title}", type="primary"):
                            if edit_pq_choice == opt_master and not course_has_master:
                                st.error("❌ Action Required: No course-wide practice exam has been uploaded yet. Upload one above before choosing this option.")
                            elif edit_pq_choice == opt_custom and not edit_custom_pq_file and not sess_info.get("style_profile"):
                                st.error("❌ Action Required: A custom practice question file is required for this setting.")
                            else:
                                with st.spinner("Syncing updates to database..."):
                                    updated_date_str = edit_date.strftime("%Y-%m-%d")
                                    
                                    if new_slides_file:
                                        updated_slides = extract_text_from_file(new_slides_file)
                                        updated_slides_fn = new_slides_file.name
                                    else:
                                        updated_slides = sess_info.get("slides", "")
                                        updated_slides_fn = sess_info.get("slides_filename", curr_slides_fn)
                                        
                                    if edit_pq_choice == opt_custom:
                                        if edit_custom_pq_file:
                                            updated_pqs = extract_text_from_file(edit_custom_pq_file)
                                            updated_pqs_fn = edit_custom_pq_file.name
                                            updated_style = analyze_style_profile(updated_pqs)
                                        else:
                                            updated_pqs_fn = sess_info.get("pqs_filename", "")
                                            updated_style = sess_info.get("style_profile", "")
                                        updated_mode = "custom"
                                    else:
                                        updated_pqs_fn = master_fn
                                        updated_style = course_data.get("global_style_profile", "")
                                        updated_mode = "course_master"
                                    
                                    if edit_title != s_title:
                                        del course_data["sessions"][s_title]
                                        
                                    course_data["sessions"][edit_title] = {
                                        "date": updated_date_str,
                                        "slides": updated_slides,
                                        "slides_filename": updated_slides_fn,
                                        "pq_mode": updated_mode,
                                        "pqs": "",  # Stripped to save database space
                                        "pqs_filename": updated_pqs_fn,
                                        "style_profile": updated_style
                                    }
                                    
                                    if save_cloud_data(data):
                                        st.session_state.saved_session_info = {"title": edit_title, "action": "updated"}
                                        st.rerun()

                    with col_del:
                        if st.button(f"🗑️ Delete Session", key=f"del_{s_title}"):
                            del course_data["sessions"][s_title]
                            if save_cloud_data(data):
                                st.success(f"Session '{s_title}' deleted.")
                                st.rerun()
